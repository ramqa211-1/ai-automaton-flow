"""
Build final presentation:
1. Add business logo to every slide
2. Embed infographic PNGs as full-slide image slides
3. Save to output path

Usage: python build_final_presentation.py input.pptx output.pptx infographic1.png [infographic2.png ...]
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.text import PP_ALIGN

LOGO_PATH = r"C:\Users\RamWalastal\ai-automaton-flow\images-buis\logo.png"
LOGO_WIDTH = Inches(1.2)
LOGO_MARGIN = Inches(0.15)


def add_logo(slide, slide_width):
    left = slide_width - LOGO_WIDTH - LOGO_MARGIN
    top = LOGO_MARGIN
    pic = slide.shapes.add_picture(LOGO_PATH, left, top, width=LOGO_WIDTH)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_infographic_slide(prs, image_path: str):
    """Add a full-slide infographic image as a new slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Fill entire slide with the image
    slide.shapes.add_picture(image_path, 0, 0, width=slide_width, height=slide_height)

    # Add logo on top
    add_logo(slide, slide_width)

    return slide


def build_final_presentation(input_path: str, output_path: str, infographic_paths: list[str]):
    prs = Presentation(input_path)
    slide_width = prs.slide_width

    # Add logo to all existing slides
    for slide in prs.slides:
        add_logo(slide, slide_width)

    # Embed each infographic as a full slide at the end
    for img_path in infographic_paths:
        if not Path(img_path).exists():
            print(f"Warning: infographic not found: {img_path}")
            continue
        add_infographic_slide(prs, img_path)
        print(f"Embedded: {img_path}")

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python build_final_presentation.py input.pptx output.pptx [infographic1.png ...]")
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2]
    infographics = sys.argv[3:]
    build_final_presentation(inp, out, infographics)
