"""
Add business logo to every slide in a PPTX file.
Usage: python add_logo_to_pptx.py input.pptx [output.pptx]
If output is omitted, overwrites the input file.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

LOGO_PATH = r"C:\Users\RamWalastal\ai-automaton-flow\images-buis\logo.png"
LOGO_WIDTH = Inches(1.2)
LOGO_MARGIN = Inches(0.15)


def add_logo_to_pptx(input_path: str, output_path: str | None = None):
    prs = Presentation(input_path)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide in prs.slides:
        # Place logo in top-right corner
        left = slide_width - LOGO_WIDTH - LOGO_MARGIN
        top = LOGO_MARGIN

        pic = slide.shapes.add_picture(LOGO_PATH, left, top, width=LOGO_WIDTH)
        # Move logo to back so it doesn't cover content
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

    out = output_path or input_path
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python add_logo_to_pptx.py input.pptx [output.pptx]")
        sys.exit(1)
    inp = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else None
    add_logo_to_pptx(inp, out)
